"""G6: exports must carry the content, not merely exist."""
import json
from pathlib import Path

import pytest
from pptx import Presentation

from sasb.build.html import load_content
from sasb.build.pptx import render
from sasb.build.svg import load_diagrams
from sasb.i18n import LOCALES, load_locale
from sasb.model import load_entities

RECON = json.loads(Path("state/reconciliation.json").read_text(encoding="utf-8"))
DIST = Path("dist")


def _texts(prs) -> str:
    out = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                out.append(shape.text_frame.text)
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    out += [cell.text for cell in row.cells]
    return "\n".join(out)


@pytest.mark.skipif(not (DIST / "reference-architecture.en@4x.png").exists(),
                    reason="run sasb.build.raster first")
@pytest.mark.parametrize("locale", list(LOCALES))
def test_pptx_carries_every_entity_and_its_provenance(tmp_path, locale):
    out = render(load_content(Path("model/content.yaml")),
                 {e.id: e for e in load_entities(Path("model/entities.yaml"))},
                 load_diagrams(Path("model/diagram.yaml")), RECON,
                 load_locale(locale), tmp_path / "b.pptx", DIST, locale)
    prs = Presentation(str(out))
    text = _texts(prs)
    assert len(prs.slides) >= 12
    assert RECON["article"]["author"] in text, "attribution lost in export"
    assert RECON["checked_at"] in text, "provenance lost in export"
    for row in RECON["entities"]:
        assert row["name"] in text, f"{row['name']} missing from the deck"
    assert "{" not in text, "unresolved entity token reached the deck"
    pictures = sum(1 for s in prs.slides for sh in s.shapes if sh.shape_type == 13)
    assert pictures >= 3, "diagrams missing from the deck"
