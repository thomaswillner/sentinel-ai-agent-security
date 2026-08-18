"""Build the deck from the canonical model with python-pptx.

Deliberately not an HTML conversion. Slides are real text frames and a real
table, so the deck stays editable in PowerPoint; diagrams are embedded as 4x
PNGs rendered from our own SVG. Pure: no network, no clock -- the timestamp
comes from the reconciliation document.
"""
from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

from ..i18n import LOCALES, load_locale
from ..model import Entity, load_entities
from .html import AVAIL_LABEL, VERDICT_LABEL, _ui, load_content, resolve_tokens
from .svg import load_diagrams

ROOT = Path(__file__).resolve().parents[3]
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
INK = RGBColor(0x14, 0x20, 0x3A)
MUTED = RGBColor(0x4A, 0x58, 0x76)
ACCENT = RGBColor(0x2A, 0x5B, 0xD7)
ROWS_PER_TABLE_SLIDE = 12
BADGE = {
    "CURRENT": RGBColor(0x1A, 0x7F, 0x4B), "CHANGED": RGBColor(0x1F, 0x5F, 0xA8),
    "RENAMED": RGBColor(0x8A, 0x5A, 0x00), "DEPRECATED": RGBColor(0x8C, 0x3A, 0x2E),
}


def _blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _textbox(slide, left, top, width, height, text, *, size=18, bold=False,
             color=INK, wrap=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = wrap
    para = frame.paragraphs[0]
    run = para.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return box


def _bullets(slide, left, top, width, height, items, *, size=16):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    for index, item in enumerate(items):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.space_after = Pt(10)
        run = para.add_run()
        run.text = item
        run.font.size = Pt(size)
        run.font.color.rgb = INK
        run.font.name = "Segoe UI"
    return box


def _title_slide(prs, loc, recon):
    slide = _blank(prs)
    art = recon["article"]
    _textbox(slide, Inches(0.9), Inches(2.0), Inches(11.5), Inches(1.4),
             _ui(loc, "brief_title"), size=40, bold=True)
    _textbox(slide, Inches(0.9), Inches(3.3), Inches(11.5), Inches(1.0),
             _ui(loc, "brief_subtitle"), size=18, color=MUTED)
    _textbox(slide, Inches(0.9), Inches(4.6), Inches(11.5), Inches(0.5),
             f'{_ui(loc, "derived_from")} "{art["subject"]}" '
             f'{_ui(loc, "by_author")} {art["author"]}, {art["post_time"][:10]}',
             size=13, color=MUTED)
    _textbox(slide, Inches(0.9), Inches(5.1), Inches(11.5), Inches(0.5),
             f'{_ui(loc, "reconciled_prefix")} {recon["checked_at"]}',
             size=13, color=ACCENT)


def _section_slide(prs, loc, sid, entities):
    section = (loc.get("sections") or {}).get(sid, {})
    slide = _blank(prs)
    _textbox(slide, Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.9),
             section.get("heading", sid), size=30, bold=True)
    items = [resolve_tokens(t, entities) for t in section.get("body", [])]
    if section.get("bullets"):
        items += [f"• {resolve_tokens(b, entities)}" for b in section["bullets"]]
    _bullets(slide, Inches(0.8), Inches(1.6), Inches(11.7), Inches(5.4), items,
             size=15 if len(items) > 2 else 17)


def _figure_slide(prs, loc, figure_id, png: Path):
    """Insert a figure, aspect-ratio fitted.

    Dimensions come from the PNG IHDR header rather than an imaging library, so
    the deck builds with no dependency beyond python-pptx.
    """
    import struct
    data = png.read_bytes()
    width_px, height_px = struct.unpack(">II", data[16:24])
    slide = _blank(prs)
    label = (loc.get("diagram") or {}).get(f"fig_{figure_id}", figure_id)
    _textbox(slide, Inches(0.6), Inches(0.3), Inches(12.1), Inches(0.7),
             label, size=24, bold=True)
    area_w, area_h = Inches(12.1), Inches(6.0)
    scale = min(area_w / width_px, area_h / height_px)
    draw_w, draw_h = Emu(int(width_px * scale)), Emu(int(height_px * scale))
    left = Emu(int((SLIDE_W - draw_w) / 2))
    slide.shapes.add_picture(str(png), left, Inches(1.2), draw_w, draw_h)


def _table_slides(prs, loc, recon):
    rows = recon["entities"]
    headers = [_ui(loc, k) for k in
               ("col_component", "col_kind", "col_status", "col_detected")]
    for start in range(0, len(rows), ROWS_PER_TABLE_SLIDE):
        chunk = rows[start:start + ROWS_PER_TABLE_SLIDE]
        slide = _blank(prs)
        _textbox(slide, Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.8),
                 _ui(loc, "section_reconciliation"), size=28, bold=True)
        shape = slide.shapes.add_table(
            len(chunk) + 1, 4, Inches(0.8), Inches(1.4), Inches(11.7),
            Inches(0.4) * (len(chunk) + 1))
        table = shape.table
        for col, header in enumerate(headers):
            cell = table.cell(0, col)
            cell.text = header
            run = cell.text_frame.paragraphs[0].runs[0]
            run.font.size, run.font.bold = Pt(12), True
            run.font.color.rgb = MUTED
        for index, row in enumerate(chunk, start=1):
            verdict = row["verdict"]
            values = [
                row["name"],
                _ui(loc, f"kind_{row['kind']}"),
                _ui(loc, VERDICT_LABEL.get(verdict, "status_current")),
                _ui(loc, AVAIL_LABEL.get(row["status_detected"], "avail_ga")),
            ]
            for col, value in enumerate(values):
                cell = table.cell(index, col)
                cell.text = value
                run = cell.text_frame.paragraphs[0].runs[0]
                run.font.size = Pt(11)
                run.font.color.rgb = BADGE.get(verdict, INK) if col == 2 else INK
                run.font.bold = col == 2


def render(content: dict, entities: dict[str, Entity], figures, recon: dict,
           loc: dict, out: Path, dist: Path, locale: str) -> Path:
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    _title_slide(prs, loc, recon)
    fig_ids = {f.id for f in figures}
    for spec in content["sections"]:
        _section_slide(prs, loc, spec["id"], entities)
        figure_id = spec.get("figure")
        if figure_id and figure_id in fig_ids:
            png = dist / f"{figure_id}.{locale}@4x.png"
            if png.exists():
                _figure_slide(prs, loc, figure_id, png)
    _table_slides(prs, loc, recon)
    out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(out))
    return out


def main(argv: list[str] | None = None) -> int:
    entities = {e.id: e for e in load_entities(ROOT / "model" / "entities.yaml")}
    figures = load_diagrams(ROOT / "model" / "diagram.yaml")
    content = load_content(ROOT / "model" / "content.yaml")
    recon = json.loads((ROOT / "state" / "reconciliation.json").read_text(encoding="utf-8"))
    dist = ROOT / "dist"
    for locale in LOCALES:
        out = render(content, entities, figures, recon, load_locale(locale),
                     dist / f"brief.{locale}.pptx", dist, locale)
        print(f"  {out.name}  {out.stat().st_size:,} bytes")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
